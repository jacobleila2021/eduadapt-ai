"""In-memory, curriculum-neutral multi-format ingestion with local-first OCR."""

from __future__ import annotations

import hashlib
import io
import re
import zipfile
from html.parser import HTMLParser
from pathlib import PurePosixPath
from typing import Any, Callable

from engines.knowledge_ingestion_engine.source_schema import (
    ContentBlock,
    SourceDocumentEnvelope,
)

MAX_SOURCE_BYTES = 50 * 1024 * 1024
MAX_EXPANDED_BYTES = 200 * 1024 * 1024
SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".txt",
    ".md",
    ".markdown",
    ".html",
    ".htm",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
}


class _TextHTMLParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self._ignored = 0

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag.lower() in {"script", "style", "noscript"}:
            self._ignored += 1

    def handle_endtag(self, tag: str) -> None:
        if tag.lower() in {"script", "style", "noscript"} and self._ignored:
            self._ignored -= 1

    def handle_data(self, data: str) -> None:
        if not self._ignored and data.strip():
            self.parts.append(data.strip())


def _error(
    stage: str,
    code: str,
    reason: str,
    recovery: str,
    *,
    retryable: bool = False,
    fallback_used: str = "none",
) -> dict[str, Any]:
    return {
        "stage": stage,
        "code": code,
        "reason": reason,
        "safe_message": reason,
        "recovery": recovery,
        "retryable": retryable,
        "fallback_used": fallback_used,
    }


def _extension(filename: str) -> str:
    name = (filename or "").lower().strip()
    return "." + name.rsplit(".", 1)[-1] if "." in name else ""


def _clean_extracted_block(text: str) -> str:
    """Repair common PDF glyph and word-boundary extraction artefacts."""
    clean = re.sub(
        r"[\x00-\x08\x0b\x0c\x0e-\x1f\ud800-\udfff\ufeff]",
        "",
        text or "",
    )
    clean = re.sub(r"/?square\d*", ". ", clean, flags=re.I)
    clean = re.sub(
        r"\b(Activity\s*\d+(?:\.\d+)*)(?:\s*\1)+",
        r"\1",
        clean,
        flags=re.I,
    )
    clean = re.sub(
        r"\b(Figure\s*\d+(?:\.\d+)*)(?:\s*\1)+",
        r"\1",
        clean,
        flags=re.I,
    )
    clean = re.sub(r"(\d+(?:\.\d+)+)(?=[A-Z][a-z])", r"\1. ", clean)
    clean = re.sub(r"(?<=[a-z)])(?=[A-Z][a-z])", ". ", clean)
    clean = re.sub(r"[ \t]+", " ", clean)
    clean = re.sub(r"\s+\n", "\n", clean)
    clean = re.sub(r"\n{3,}", "\n\n", clean)
    return clean.strip()


def _detect_mime(extension: str, data: bytes) -> tuple[str, bool]:
    signatures = {
        ".pdf": ("application/pdf", data.startswith(b"%PDF")),
        ".png": ("image/png", data.startswith(b"\x89PNG\r\n\x1a\n")),
        ".jpg": ("image/jpeg", data.startswith(b"\xff\xd8\xff")),
        ".jpeg": ("image/jpeg", data.startswith(b"\xff\xd8\xff")),
        ".webp": (
            "image/webp",
            len(data) >= 12 and data[:4] == b"RIFF" and data[8:12] == b"WEBP",
        ),
        ".docx": ("application/vnd.openxmlformats-officedocument.wordprocessingml.document", data.startswith(b"PK")),
        ".pptx": ("application/vnd.openxmlformats-officedocument.presentationml.presentation", data.startswith(b"PK")),
    }
    if extension in signatures:
        return signatures[extension]
    if extension in {".txt", ".md", ".markdown"}:
        return "text/plain", True
    if extension in {".html", ".htm"}:
        return "text/html", True
    return "application/octet-stream", False


def _validate_office_zip(data: bytes) -> str | None:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            expanded = 0
            for info in archive.infolist():
                path = PurePosixPath(info.filename.replace("\\", "/"))
                if path.is_absolute() or ".." in path.parts:
                    return "Office archive contains an unsafe member path"
                expanded += max(info.file_size, 0)
                if expanded > MAX_EXPANDED_BYTES:
                    return "Office archive expands beyond the allowed size"
                if info.compress_size and info.file_size / info.compress_size > 200:
                    return "Office archive has an unsafe compression ratio"
    except (zipfile.BadZipFile, OSError):
        return "Office document archive is corrupt"
    return None


def _ocr_image(
    image,
    *,
    cloud_ocr: Callable[[bytes], dict[str, Any]] | None = None,
    raw_bytes: bytes | None = None,
) -> tuple[str, float | None, str, list[dict[str, Any]]]:
    warnings: list[dict[str, Any]] = []
    try:
        import pytesseract
        from pytesseract import Output

        data = pytesseract.image_to_data(image, output_type=Output.DICT)
        words: list[str] = []
        confidences: list[float] = []
        for word, confidence in zip(data.get("text", []), data.get("conf", [])):
            if str(word).strip():
                words.append(str(word).strip())
                try:
                    value = float(confidence)
                    if value >= 0:
                        confidences.append(value)
                except (TypeError, ValueError):
                    continue
        text = " ".join(words)
        score = round(sum(confidences) / len(confidences) / 100, 3) if confidences else None
        if text.strip() and (score is None or score >= 0.55):
            return text, score, "tesseract", warnings
        warnings.append(
            _error(
                "ocr",
                "low_ocr_confidence",
                "Local OCR confidence is too low for reliable adaptation.",
                "Upload a clearer scan or a text-based document.",
                retryable=True,
                fallback_used="local_ocr",
            )
        )
    except Exception:
        warnings.append(
            _error(
                "ocr",
                "local_ocr_unavailable",
                "Local OCR is unavailable on this deployment.",
                "Install Tesseract OCR or upload a text-based document.",
                retryable=False,
            )
        )

    if cloud_ocr and raw_bytes:
        try:
            result = cloud_ocr(raw_bytes) or {}
            text = str(result.get("text") or "")
            confidence = result.get("confidence")
            if text.strip() and (confidence is None or float(confidence) >= 0.55):
                return text, float(confidence) if confidence is not None else None, "cloud_vision", warnings
        except Exception:
            warnings.append(
                _error(
                    "ocr",
                    "cloud_ocr_failed",
                    "The optional OCR fallback could not read this image.",
                    "Upload a clearer scan or a text-based document.",
                    retryable=True,
                    fallback_used="cloud_vision",
                )
            )
    return "", None, "none", warnings


def _language_hint(text: str) -> str:
    sample = f" {text[:5000].lower()} "
    french = sum(sample.count(word) for word in (" le ", " la ", " les ", " une ", " des ", " et "))
    english = sum(sample.count(word) for word in (" the ", " is ", " are ", " and ", " of ", " to "))
    if french >= 3 and french > english:
        return "fr"
    if english >= 2 or re.search(r"[A-Za-z]{4}", sample):
        return "en"
    return "unknown"


def _make_envelope(
    filename: str,
    data: bytes,
    mime: str,
    extension: str,
    blocks: list[ContentBlock],
    warnings: list[dict[str, Any]],
    errors: list[dict[str, Any]],
) -> SourceDocumentEnvelope:
    text = "\n\n".join(block.text.strip() for block in blocks if block.text.strip())
    source_hash = hashlib.sha256(data).hexdigest()
    methods = sorted({block.extraction_method for block in blocks})
    ocr_scores = [
        block.ocr_confidence
        for block in blocks
        if block.ocr_confidence is not None
    ]
    printable = len(re.findall(r"\w", text, re.UNICODE))
    score = min(1.0, printable / 120.0) if text else 0.0
    if errors and not text:
        status = "rejected"
    elif score < 0.15:
        status = "unreadable"
    elif warnings:
        status = "partial"
    else:
        status = "readable"
    return SourceDocumentEnvelope(
        schema_version="3.0.0",
        source_id=f"src_{source_hash[:16]}",
        source_hash=source_hash,
        filename=filename,
        detected_mime=mime,
        detected_format=extension.lstrip("."),
        status=status,
        blocks=blocks,
        text=text,
        extraction_methods=methods,
        ocr_used=any(method in {"tesseract", "cloud_vision"} for method in methods),
        ocr_confidence=(
            round(sum(ocr_scores) / len(ocr_scores), 3) if ocr_scores else None
        ),
        readable_content_score=round(score, 3),
        language=_language_hint(text),
        warnings=warnings,
        errors=errors,
    )


def ingest_source_bytes(
    filename: str,
    file_bytes: bytes,
    *,
    cloud_ocr: Callable[[bytes], dict[str, Any]] | None = None,
    user_metadata: dict[str, Any] | None = None,
) -> SourceDocumentEnvelope:
    """Extract an uploaded educational source without assuming any curriculum."""
    data = bytes(file_bytes or b"")
    extension = _extension(filename)
    source_hash = hashlib.sha256(data).hexdigest()
    mime, signature_ok = _detect_mime(extension, data)
    if extension not in SUPPORTED_EXTENSIONS:
        return SourceDocumentEnvelope(
            schema_version="3.0.0",
            source_id=f"src_{source_hash[:16]}",
            source_hash=source_hash,
            filename=filename,
            detected_mime=mime,
            detected_format=extension.lstrip("."),
            status="rejected",
            errors=[
                _error(
                    "validation",
                    "unsupported_format",
                    "This file format is not supported.",
                    "Upload PDF, DOCX, PPTX, TXT, Markdown, HTML, PNG, JPEG, or WebP.",
                )
            ],
        )
    if not data:
        return SourceDocumentEnvelope(
            schema_version="3.0.0",
            source_id=f"src_{source_hash[:16]}",
            source_hash=source_hash,
            filename=filename,
            detected_mime=mime,
            detected_format=extension.lstrip("."),
            status="rejected",
            errors=[
                _error(
                    "validation",
                    "empty_input",
                    "The uploaded file is empty.",
                    "Choose a file containing readable educational content.",
                )
            ],
        )
    if len(data) > MAX_SOURCE_BYTES:
        return SourceDocumentEnvelope(
            schema_version="3.0.0",
            source_id=f"src_{source_hash[:16]}",
            source_hash=source_hash,
            filename=filename,
            detected_mime=mime,
            detected_format=extension.lstrip("."),
            status="rejected",
            errors=[
                _error(
                    "validation",
                    "file_too_large",
                    "The uploaded file exceeds the 50 MB limit.",
                    "Split the document into smaller files.",
                )
            ],
        )
    if not signature_ok:
        return SourceDocumentEnvelope(
            schema_version="3.0.0",
            source_id=f"src_{source_hash[:16]}",
            source_hash=source_hash,
            filename=filename,
            detected_mime=mime,
            detected_format=extension.lstrip("."),
            status="rejected",
            errors=[
                _error(
                    "validation",
                    "signature_mismatch",
                    "The file content does not match its extension.",
                    "Export the original document again and re-upload it.",
                )
            ],
        )
    if extension in {".pdf", ".png", ".jpg", ".jpeg", ".webp"} and b"PK\x03\x04" in data[:8192]:
        return SourceDocumentEnvelope(
            schema_version="3.0.0",
            source_id=f"src_{source_hash[:16]}",
            source_hash=source_hash,
            filename=filename,
            detected_mime=mime,
            detected_format=extension.lstrip("."),
            status="rejected",
            errors=[
                _error(
                    "validation",
                    "polyglot_file_rejected",
                    "The file contains conflicting format signatures.",
                    "Export a clean copy from the original application.",
                )
            ],
        )
    if extension in {".docx", ".pptx"}:
        unsafe = _validate_office_zip(data)
        if unsafe:
            return SourceDocumentEnvelope(
                schema_version="3.0.0",
                source_id=f"src_{source_hash[:16]}",
                source_hash=source_hash,
                filename=filename,
                detected_mime=mime,
                detected_format=extension.lstrip("."),
                status="rejected",
                errors=[
                    _error(
                        "validation",
                        "unsafe_office_archive",
                        unsafe,
                        "Export a fresh DOCX or PPTX and upload it again.",
                    )
                ],
            )

    blocks: list[ContentBlock] = []
    warnings: list[dict[str, Any]] = []
    errors: list[dict[str, Any]] = []

    def add(
        text: str,
        *,
        kind: str = "prose",
        page: int | None = None,
        slide: int | None = None,
        method: str = "native",
        confidence: float | None = None,
        metadata: dict[str, Any] | None = None,
    ) -> None:
        clean = _clean_extracted_block(text)
        if not clean:
            return
        order = len(blocks)
        blocks.append(
            ContentBlock(
                block_id=f"blk_{source_hash[:10]}_{order:05d}",
                text=clean,
                kind=kind,  # type: ignore[arg-type]
                page=page,
                slide=slide,
                order=order,
                extraction_method=method,
                ocr_confidence=confidence,
                source_ref=f"{filename}#{'slide' if slide else 'page'}-{slide or page or 1}",
                metadata=metadata or {},
            )
        )

    try:
        if extension == ".pdf":
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(data))
            fitz_document = None
            for index, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ""
                if len(re.findall(r"\w", text)) >= 12:
                    add(text, page=index, method="pypdf")
                    continue
                try:
                    import fitz
                    from PIL import Image

                    if fitz_document is None:
                        fitz_document = fitz.open(stream=data, filetype="pdf")
                    pixmap = fitz_document[index - 1].get_pixmap(
                        matrix=fitz.Matrix(2, 2)
                    )
                    raw = pixmap.tobytes("png")
                    image = Image.open(io.BytesIO(raw))
                    ocr_text, confidence, method, ocr_warnings = _ocr_image(
                        image, cloud_ocr=cloud_ocr, raw_bytes=raw
                    )
                    warnings.extend(ocr_warnings)
                    add(
                        ocr_text,
                        page=index,
                        kind="image_text",
                        method=method,
                        confidence=confidence,
                    )
                except Exception:
                    warnings.append(
                        _error(
                            "ocr",
                            "scanned_page_unreadable",
                            f"Page {index} contains no reliably extractable text.",
                            "Upload a clearer scan or run OCR before uploading.",
                            retryable=True,
                        )
                    )
            if fitz_document is not None:
                fitz_document.close()
        elif extension == ".docx":
            from docx import Document

            document = Document(io.BytesIO(data))
            for paragraph in document.paragraphs:
                kind = "heading" if (paragraph.style and str(paragraph.style.name).lower().startswith("heading")) else "prose"
                add(paragraph.text, kind=kind, method="python-docx")
            for table_index, table in enumerate(document.tables, 1):
                rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
                add(
                    "\n".join(rows),
                    kind="table",
                    method="python-docx",
                    metadata={"table_index": table_index},
                )
            try:
                from PIL import Image

                with zipfile.ZipFile(io.BytesIO(data)) as archive:
                    media = [
                        name
                        for name in archive.namelist()
                        if name.lower().startswith("word/media/")
                    ]
                    for name in media[:40]:
                        raw = archive.read(name)
                        image = Image.open(io.BytesIO(raw))
                        text, confidence, method, ocr_warnings = _ocr_image(
                            image, cloud_ocr=cloud_ocr, raw_bytes=raw
                        )
                        warnings.extend(ocr_warnings)
                        add(
                            text,
                            kind="image_text",
                            method=method,
                            confidence=confidence,
                            metadata={"embedded_asset": name},
                        )
            except Exception:
                pass
        elif extension == ".pptx":
            try:
                from pptx import Presentation
            except ImportError as exc:
                raise RuntimeError("PowerPoint support is unavailable") from exc
            presentation = Presentation(io.BytesIO(data))
            for slide_index, slide in enumerate(presentation.slides, 1):
                for shape in slide.shapes:
                    text = str(getattr(shape, "text", "") or "")
                    add(text, slide=slide_index, method="python-pptx")
                    image = getattr(shape, "image", None)
                    if image is not None:
                        try:
                            from PIL import Image

                            raw = image.blob
                            ocr_text, confidence, method, ocr_warnings = _ocr_image(
                                Image.open(io.BytesIO(raw)),
                                cloud_ocr=cloud_ocr,
                                raw_bytes=raw,
                            )
                            warnings.extend(ocr_warnings)
                            add(
                                ocr_text,
                                kind="image_text",
                                slide=slide_index,
                                method=method,
                                confidence=confidence,
                            )
                        except Exception:
                            pass
                try:
                    notes = slide.notes_slide.notes_text_frame.text
                    add(
                        notes,
                        kind="speaker_notes",
                        slide=slide_index,
                        method="python-pptx",
                    )
                except Exception:
                    pass
        elif extension in {".png", ".jpg", ".jpeg", ".webp"}:
            from PIL import Image

            text, confidence, method, ocr_warnings = _ocr_image(
                Image.open(io.BytesIO(data)),
                cloud_ocr=cloud_ocr,
                raw_bytes=data,
            )
            warnings.extend(ocr_warnings)
            add(
                text,
                kind="image_text",
                page=1,
                method=method,
                confidence=confidence,
            )
        else:
            decoded = data.decode("utf-8-sig", errors="replace")
            if extension in {".html", ".htm"}:
                parser = _TextHTMLParser()
                parser.feed(decoded)
                decoded = "\n".join(parser.parts)
            add(decoded, page=1, method="text")
    except Exception as exc:
        errors.append(
            _error(
                "extraction",
                "corrupt_or_unreadable",
                "The document could not be read safely.",
                "Export a fresh copy or upload plain text.",
                retryable=False,
            )
        )
        warnings.append(
            {
                "stage": "extraction",
                "code": "private_diagnostic",
                "reason": type(exc).__name__,
                "safe_message": "Extraction failed.",
                "recovery": "Review private server logs.",
                "retryable": False,
                "fallback_used": "none",
            }
        )

    envelope = _make_envelope(
        filename, data, mime, extension, blocks, warnings, errors
    )
    envelope.user_metadata = dict(user_metadata or {})
    if not envelope.ok and not envelope.errors:
        envelope.errors.append(
            _error(
                "extraction",
                "unreadable_content",
                "No reliable educational text could be extracted.",
                "Upload a clearer scan, an OCR-processed file, or plain text.",
                retryable=True,
                fallback_used="local_ocr" if envelope.ocr_used else "native_extraction",
            )
        )
    return envelope
